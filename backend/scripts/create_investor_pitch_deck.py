import sys
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    BG_DARK = RGBColor(15, 23, 42)      # Slate 900
    CARD_BG = RGBColor(30, 41, 59)      # Slate 800
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(148, 163, 184) # Slate 400
    ORANGE_ACCENT = RGBColor(249, 115, 22) # Orange 500
    BLUE_ACCENT = RGBColor(99, 102, 241)   # Indigo 500
    GREEN_ACCENT = RGBColor(16, 185, 129)  # Emerald 500

    def add_header(slide, title_text, category_text="GO2PICK INVESTOR PITCH"):
        # Header category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ORANGE_ACCENT

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

    def set_slide_background(slide, color=BG_DARK):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    # ─── SLIDE 1: Title Slide ───────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, BG_DARK)

    # Decorative Card Box
    shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = ORANGE_ACCENT
    shape.line.width = Pt(2)

    tb = slide1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.9), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "Go2Pick"
    p0.font.size = Pt(48)
    p0.font.bold = True
    p0.font.color.rgb = ORANGE_ACCENT
    p0.alignment = PP_ALIGN.LEFT

    p1 = tf.add_paragraph()
    p1.text = "Revolutionizing Hyperlocal Commerce via Zero-Wait Pre-Order & Pickup"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE

    p2 = tf.add_paragraph()
    p2.text = "\nEliminating 25% delivery commissions & long checkout queues for local stores."
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_MUTED

    p3 = tf.add_paragraph()
    p3.text = "\n\nPresented by: Go2Pick Founding Team\nContact: founder@go2pick.com | Live Demo: go2-pick.vercel.app"
    p3.font.size = Pt(13)
    p3.font.color.rgb = GREEN_ACCENT

    # ─── SLIDE 2: The Problem ────────────────────────────────────────────────
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, BG_DARK)
    add_header(slide2, "The Problem: Broken Economics in Local Retail & Quick Commerce")

    problems = [
        ("High Delivery Commissions", "Platforms charge local shopkeepers 20-30% commissions, destroying small business profit margins."),
        ("In-Store Queue Fatigue", "Customers spend 15-25 minutes waiting in checkout queues during peak evening hours for daily essentials."),
        ("Delivery Fleet Burn Rate", "Quick-commerce delivery apps lose $1.50 per order on rider payouts, traffic delays, and dark store overhead.")
    ]

    for i, (title, desc) in enumerate(problems):
        left = Inches(0.8 + i * 3.9)
        top = Inches(1.8)
        width = Inches(3.7)
        height = Inches(4.8)

        box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = RGBColor(51, 65, 85)

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"0{i+1}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = ORANGE_ACCENT

        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE

        p_d = tf.add_paragraph()
        p_d.text = f"\n{desc}"
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = TEXT_MUTED

    # ─── SLIDE 3: The Solution ───────────────────────────────────────────────
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, BG_DARK)
    add_header(slide3, "The Solution: Go2Pick Pre-Order & Express Pickup Passcode")

    sol_box = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
    sol_box.fill.solid()
    sol_box.fill.fore_color.rgb = CARD_BG
    sol_box.line.color.rgb = GREEN_ACCENT

    tf3 = sol_box.text_frame
    tf3.word_wrap = True

    p = tf3.paragraphs[0]
    p.text = "How Go2Pick Works:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT

    steps = [
        "1. Customer Browses & Orders: Customers discover neighborhood stores (e.g., Grany Groceries), pick fresh items, and pre-order.",
        "2. Merchant Packing Alert: Shopkeeper receives instant alert, packs the items, and generates a unique 4-digit Pickup Passcode.",
        "3. Zero-Wait Pickup: Customer walks in, shows passcode, grabs packed bag, and exits in under 30 seconds."
    ]

    for s in steps:
        p_s = tf3.add_paragraph()
        p_s.text = f"\n• {s}"
        p_s.font.size = Pt(15)
        p_s.font.color.rgb = TEXT_WHITE

    # ─── SLIDE 4: Market Opportunity ─────────────────────────────────────────
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, BG_DARK)
    add_header(slide4, "Market Opportunity: $50 Billion Hyperlocal Retail Market")

    stats = [
        ("$50 Billion", "Total Addressable Market (TAM) for local retail & grocery in India"),
        ("$12 Billion", "Serviceable Addressable Market (SAM) in Tier-1 & Tier-2 cities"),
        ("15 Million+", "Local Kirana & Independent stores needing digital pickup enablement")
    ]

    for i, (val, desc) in enumerate(stats):
        left = Inches(0.8 + i * 3.9)
        box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.7), Inches(4.8))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = BLUE_ACCENT

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = ORANGE_ACCENT

        p_d = tf.add_paragraph()
        p_d.text = f"\n\n{desc}"
        p_d.font.size = Pt(15)
        p_d.font.color.rgb = TEXT_WHITE

    # ─── SLIDE 5: Technology Architecture ────────────────────────────────────
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, BG_DARK)
    add_header(slide5, "Technical Architecture: Scalable Real-time Marketplace")

    tech_box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
    tech_box.fill.solid()
    tech_box.fill.fore_color.rgb = CARD_BG
    tech_box.line.color.rgb = BLUE_ACCENT

    tf5 = tech_box.text_frame
    tf5.word_wrap = True

    p = tf5.paragraphs[0]
    p.text = "Built for Enterprise Reliability & Sub-100ms Performance:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT

    techs = [
        "Backend Stack: FastAPI (Python 3.11), PyMongo, Uvicorn ASGI async web server.",
        "Frontend Stack: React 18, Vite build system, Tailwind CSS & Glassmorphism UI.",
        "Database Layer: Dual persistence — Firebase Cloud Firestore & MongoDB Atlas.",
        "Auth & Push: Firebase Authentication, FCM Push Notifications, JWT Tokens.",
        "Deployment: Vercel Serverless Edge deployment for global sub-50ms latency."
    ]

    for t in techs:
        pt = tf5.add_paragraph()
        pt.text = f"\n⚡ {t}"
        pt.font.size = Pt(15)
        pt.font.color.rgb = TEXT_WHITE

    # ─── SLIDE 6: Business Model ─────────────────────────────────────────────
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, BG_DARK)
    add_header(slide6, "Monetization Model: Multi-Stream Sustainable Revenue")

    models = [
        ("Merchant SaaS Subscription", "₹499 - ₹1,499 / month per store for inventory management & analytics."),
        ("Featured Shop Placements", "Sponsored top placement in Customer Home & Search Explore for local shops."),
        ("Convenience Pickup Fee", "Nominal ₹5 - ₹10 convenience packaging fee paid by customer per pre-order.")
    ]

    for i, (title, desc) in enumerate(models):
        left = Inches(0.8 + i * 3.9)
        box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.7), Inches(4.8))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = GREEN_ACCENT

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = GREEN_ACCENT

        p_d = tf.add_paragraph()
        p_d.text = f"\n\n{desc}"
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = TEXT_WHITE

    # ─── SLIDE 7: Financial Projections ──────────────────────────────────────
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, BG_DARK)
    add_header(slide7, "Financial Projections: 3-Year Growth Roadmap")

    box7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
    box7.fill.solid()
    box7.fill.fore_color.rgb = CARD_BG
    box7.line.color.rgb = ORANGE_ACCENT

    tf7 = box7.text_frame
    tf7.word_wrap = True

    p = tf7.paragraphs[0]
    p.text = "Year 1 - Year 3 Projected Milestones:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE_ACCENT

    proj = [
        "Year 1: 500 Active Shops | 50,000 Monthly Orders | $120,000 ARR",
        "Year 2: 3,500 Active Shops | 450,000 Monthly Orders | $1.2 Million ARR",
        "Year 3: 15,000 Active Shops | 2.5 Million Monthly Orders | $6.5 Million ARR"
    ]

    for pr in proj:
        pp = tf7.add_paragraph()
        pp.text = f"\n🚀 {pr}"
        pp.font.size = Pt(16)
        pp.font.bold = True
        pp.font.color.rgb = TEXT_WHITE

    # ─── SLIDE 8: The Ask ────────────────────────────────────────────────────
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, BG_DARK)
    add_header(slide8, "The Ask: Raising $500,000 Seed Funding")

    ask_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
    ask_box.fill.solid()
    ask_box.fill.fore_color.rgb = CARD_BG
    ask_box.line.color.rgb = GREEN_ACCENT

    tf8 = ask_box.text_frame
    tf8.word_wrap = True

    p = tf8.paragraphs[0]
    p.text = "Seed Round Allocation ($500,000 for 18-Month Runway):"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT

    allocs = [
        "40% - Merchant Acquisition & Field Operations Team (5 City Hubs)",
        "35% - Technology & Product Development (AI Demand Forecasting & POS Integration)",
        "15% - Customer Hyperlocal Marketing & Referral Campaigns",
        "10% - Legal, Compliance & Reserve Runway"
    ]

    for a in allocs:
        pa = tf8.add_paragraph()
        pa.text = f"\n💼 {a}"
        pa.font.size = Pt(15)
        pa.font.color.rgb = TEXT_WHITE

    # Save presentation
    output_path = Path(__file__).resolve().parent.parent.parent / "Go2Pick_Investor_Pitch_Deck.pptx"
    prs.save(str(output_path))
    print(f"\n[SUCCESS] Created Investor Pitch Deck PowerPoint at: {output_path}")

if __name__ == "__main__":
    create_deck()
